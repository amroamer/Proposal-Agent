"""Phase 1 — section_splitter tests.

Three required fixtures (per spec):
  1. Full template (all 13 canonical sections present, AR titles)
  2. Missing optional section (one canonical section omitted)
  3. Mixed AR/EN headings (some markers in English)

Plus a regression test that table-cell text on CV/references slides
is captured — without this, team_structure and detailed_experience
extraction would lose its primary evidence.
"""
from __future__ import annotations

import io

import pytest
from pptx import Presentation

from app.services.proposal_review.section_mapping import SECTION_KEYS
from app.services.proposal_review.section_splitter import (
    ProposalSections,
    SlideText,
    split_pptx,
    split_slide_texts,
)


# ---- fixture builders --------------------------------------------------------

# Canonical AR titles in template order — these are the strings that
# appear after the "NN – " marker on each chapter divider slide.
_AR_TITLES = (
    "معايير التقييم",
    "الملخص التنفيذي",
    "فهمنا لمتطلباتكم",
    "القيمة التي نقدمها",
    "وجهة نظرنا",
    "نهجنا المُفصّل والجدول الزمني",
    "الهيكل التنظيمي والسير الذاتية",
    "الخبرات التفصيلية في مشاريع مشابهة",
    "الأدوات والمنهجيات",
    "الملف التعريفي بشركة كي بي إم جي",
    "السجلات والشهادات النظامية",
    "الشروط والأحكام العامة",
    "الافتراضات",
)


def _full_template_slides() -> list[SlideText]:
    """13 marker slides + a body slide between each. Slide 1 = cover."""
    slides: list[SlideText] = [
        SlideText(slide_number=1, text="KPMG Technical Proposal\nClient: Ministry of Example"),
    ]
    n = 2
    for i, title in enumerate(_AR_TITLES, start=1):
        slides.append(SlideText(slide_number=n, text=f"{i:02d} – {title}"))
        n += 1
        slides.append(SlideText(slide_number=n, text=f"Body of section {i:02d}"))
        n += 1
    return slides


def _missing_section_slides() -> list[SlideText]:
    """Same as full template but skip section 09 (tools_methodologies)."""
    slides: list[SlideText] = [
        SlideText(slide_number=1, text="Cover slide"),
    ]
    n = 2
    for i, title in enumerate(_AR_TITLES, start=1):
        if i == 9:
            continue  # skip Tools & Methodologies entirely
        slides.append(SlideText(slide_number=n, text=f"{i:02d} – {title}"))
        n += 1
        slides.append(SlideText(slide_number=n, text=f"Body of section {i:02d}"))
        n += 1
    return slides


def _mixed_lang_slides() -> list[SlideText]:
    """Half the markers in English, half in Arabic. KPMG decks built by
    bilingual teams in a hurry sometimes commit half-translated TOCs."""
    en_titles = {
        2: "Executive Summary",
        4: "Value We Provide",
        7: "Org Structure & CVs",
        9: "Tools & Methodologies",
    }
    slides: list[SlideText] = [
        SlideText(slide_number=1, text="Cover")
    ]
    n = 2
    for i, ar_title in enumerate(_AR_TITLES, start=1):
        title = en_titles.get(i, ar_title)
        slides.append(SlideText(slide_number=n, text=f"{i:02d} – {title}"))
        n += 1
        slides.append(SlideText(slide_number=n, text="body"))
        n += 1
    return slides


# ---- tests ------------------------------------------------------------------


class TestSplitFullTemplate:
    def test_finds_all_13_canonical_sections(self):
        result = split_slide_texts(_full_template_slides())
        for key in SECTION_KEYS:
            assert key in result.sections, f"missing canonical section {key}"
        assert result.missing_sections == []
        assert result.unknown_markers == []

    def test_section_starts_recorded_in_template_order(self):
        result = split_slide_texts(_full_template_slides())
        starts = [result.section_starts[k] for k in SECTION_KEYS]
        assert starts == sorted(starts), "section starts should be monotonically increasing"

    def test_front_matter_captures_pre_marker_slides(self):
        result = split_slide_texts(_full_template_slides())
        # Cover slide is before the first marker — must end up in front_matter
        front = result.sections.get("front_matter")
        assert front is not None
        assert any("kpmg technical proposal" in s.text.lower() for s in front.slides)

    def test_each_section_has_marker_and_body_slides(self):
        result = split_slide_texts(_full_template_slides())
        for key in SECTION_KEYS:
            sec = result.sections[key]
            # Marker slide + body slide = 2 minimum
            assert len(sec.slides) >= 2, f"{key} missing slides: got {len(sec.slides)}"
            assert sec.raw_text.strip() != ""


class TestSplitMissingOptionalSection:
    def test_missing_section_logged_not_raised(self):
        # MUST NOT raise — spec: "Missing sections: log a warning, continue."
        result = split_slide_texts(_missing_section_slides())
        assert isinstance(result, ProposalSections)

    def test_missing_section_listed_in_output(self):
        result = split_slide_texts(_missing_section_slides())
        assert "tools_methodologies" in result.missing_sections
        # Other sections still present
        assert "executive_summary" in result.sections
        assert "detailed_experience" in result.sections

    def test_other_sections_unaffected(self):
        result = split_slide_texts(_missing_section_slides())
        # 12 of 13 canonical present
        present = [k for k in SECTION_KEYS if k in result.sections]
        assert len(present) == 12


class TestSplitMixedLanguage:
    def test_english_marker_titles_recognised(self):
        result = split_slide_texts(_mixed_lang_slides())
        for key in (
            "executive_summary",
            "value_proposition",
            "team_structure",
            "tools_methodologies",
        ):
            assert key in result.sections, f"EN-titled marker missed: {key}"

    def test_arabic_markers_still_recognised_alongside_english(self):
        result = split_slide_texts(_mixed_lang_slides())
        for key in (
            "evaluation_criteria",
            "our_understanding",
            "our_perspective",
            "detailed_approach",
            "detailed_experience",
            "kpmg_profile",
            "certifications",
            "terms",
            "assumptions",
        ):
            assert key in result.sections, f"AR-titled marker missed: {key}"


class TestMarkerRegex:
    @pytest.mark.parametrize("sep", ["-", "–", "—", " - ", " – ", " — "])
    def test_separator_variants(self, sep):
        slides = [
            SlideText(slide_number=1, text="Cover"),
            SlideText(slide_number=2, text=f"02{sep}الملخص التنفيذي"),
            SlideText(slide_number=3, text="body"),
        ]
        result = split_slide_texts(slides)
        assert "executive_summary" in result.sections

    def test_marker_must_be_first_nonblank_line(self):
        # Body slide that mentions "01 – Foo" in a bullet must NOT
        # start a new section.
        slides = [
            SlideText(slide_number=1, text="Cover"),
            SlideText(slide_number=2, text="02 – الملخص التنفيذي"),
            SlideText(slide_number=3, text="some intro\n01 – this is a bullet, not a section"),
        ]
        result = split_slide_texts(slides)
        # We should have exactly one canonical section detected.
        canonical = [k for k in SECTION_KEYS if k in result.sections]
        assert canonical == ["executive_summary"]

    def test_unknown_marker_kept_as_warning(self):
        slides = [
            SlideText(slide_number=1, text="Cover"),
            SlideText(slide_number=2, text="01 – معايير التقييم"),
            SlideText(slide_number=3, text="body"),
            SlideText(slide_number=4, text="07 – Some Custom Section That Isn't In The Template"),
            SlideText(slide_number=5, text="custom body"),
        ]
        result = split_slide_texts(slides)
        assert "evaluation_criteria" in result.sections
        assert len(result.unknown_markers) == 1
        assert result.unknown_markers[0]["slide_number"] == 4


# ---- python-pptx round-trip + table capture regression ----------------------


def _build_pptx_with_table_on_marker_slide() -> bytes:
    """Build a real .pptx where the team_structure section contains a
    table (CV summary). This verifies that table-cell text is captured;
    without it, team_structure extraction in Phase 3 has no evidence."""
    prs = Presentation()
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[1]

    # Slide 1 — cover
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Cover"

    # Slide 2 — section 07 marker (Org Structure & CVs)
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "07 – الهيكل التنظيمي والسير الذاتية"

    # Slide 3 — same section, with a CV-style table
    s3 = prs.slides.add_slide(blank)
    rows, cols = 3, 3
    left = top = 914400  # 1 inch in EMUs
    width = 5_000_000
    height = 2_000_000
    table_shape = s3.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    tbl.cell(0, 0).text = "Name"
    tbl.cell(0, 1).text = "Role"
    tbl.cell(0, 2).text = "Years"
    tbl.cell(1, 0).text = "Sara Ahmed"
    tbl.cell(1, 1).text = "Engagement Partner"
    tbl.cell(1, 2).text = "18"
    tbl.cell(2, 0).text = "Khalid Saud"
    tbl.cell(2, 1).text = "Project Manager"
    tbl.cell(2, 2).text = "12"

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


class TestRealPptxTableCapture:
    def test_table_cell_text_lands_in_team_structure(self):
        data = _build_pptx_with_table_on_marker_slide()
        result = split_pptx(data)
        assert "team_structure" in result.sections
        text = result.sections["team_structure"].raw_text
        # Names + roles + headers must all appear
        assert "Sara Ahmed" in text
        assert "Engagement Partner" in text
        assert "Khalid Saud" in text
        assert "Project Manager" in text
        assert "Years" in text
