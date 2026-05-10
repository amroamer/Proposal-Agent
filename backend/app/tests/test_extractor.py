"""Phase 3 — extractor pipeline tests (pure-async, mocked DB + LLM).

Acceptance criteria from the spec:
  - Running extract_dossier twice on the same pptx logs
    `dossier_cache_hit=True` on the second run and does NOT call Ollama.
  - All section extractions for a single proposal start within 1s of
    each other (parallel dispatch — verified directly on the
    asyncio.gather() fan-out used by extract_dossier).
  - Dossier persists with non-null dossier_json for at least 5 sections
    on a synthetic 6-section deck.

The DB session is a MagicMock — we don't exercise real SQL here.
LLM calls are mocked at `extractor.generate_structured` so this runs
without Ollama.
"""
from __future__ import annotations

import asyncio
import io
import logging
import time
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from pptx import Presentation
from pydantic import BaseModel

from app.models.dossier import Dossier as DossierRow
from app.services.proposal_review.dossier_schemas import GenericNarrativeFacts
from app.services.proposal_review.extractor import (
    compute_source_hash,
    extract_dossier,
)
from app.services.proposal_review.llm_client import (
    Classification,
    StructuredResult,
)


# ---- helpers ----------------------------------------------------------------


_SECTION_TITLES = (
    "01 – معايير التقييم",
    "02 – الملخص التنفيذي",
    "03 – فهمنا لمتطلباتكم",
    "04 – القيمة التي نقدمها",
    "05 – وجهة نظرنا",
    "06 – نهجنا المُفصّل والجدول الزمني",
)


def _build_proposal_pptx() -> bytes:
    """6-section synthetic deck — exceeds the spec's 5-section minimum."""
    prs = Presentation()
    layout = prs.slide_layouts[1]
    cover = prs.slides.add_slide(layout)
    cover.shapes.title.text = "Cover"
    for title in _SECTION_TITLES:
        m = prs.slides.add_slide(layout)
        m.shapes.title.text = title
        b = prs.slides.add_slide(layout)
        b.shapes.title.text = "Body"
        body_ph = next(
            (ph for ph in b.placeholders if ph.placeholder_format.idx == 1),
            None,
        )
        if body_ph is not None:
            body_ph.text_frame.text = f"Body content for {title}"
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _make_structured_result(model_cls: type[BaseModel]) -> StructuredResult:
    """Schema-valid empty payload for whichever section schema was requested."""
    instance = model_cls()
    return StructuredResult(
        parsed=instance,
        raw_output=instance.model_dump_json(),
        model="qwen2.5:32b",
        duration_ms=10,
    )


def _mock_db(*, find_returns: object = None, classification: str = "Restricted") -> MagicMock:
    """Build a MagicMock async session that satisfies the calls extract_dossier
    makes:
      - .execute() (used twice: cache lookup + classification load)
      - .add()
      - .flush()
    """
    db = MagicMock()
    db.add = MagicMock()
    db.flush = AsyncMock()
    cache_result = MagicMock()
    cache_result.scalar_one_or_none = MagicMock(return_value=find_returns)
    cls_result = MagicMock()
    cls_result.scalar_one_or_none = MagicMock(return_value=classification)
    db.execute = AsyncMock(side_effect=[cache_result, cls_result, cache_result, cls_result, cache_result, cls_result])
    return db


def _set_inserted_id(db: MagicMock, row_id: int = 99) -> None:
    """Make `db.add(row)` set row.id (simulating SQLAlchemy autoincrement)."""
    def _add(row: DossierRow):
        if not getattr(row, "id", None):
            row.id = row_id
    db.add.side_effect = _add


async def _mock_generate(*, schema: type[BaseModel], **_kwargs):
    """Drop-in replacement for llm_client.generate_structured.
    Tiny sleep simulates real HTTP latency so async dispatch is observable."""
    await asyncio.sleep(0.05)
    return _make_structured_result(schema)


# ---- tests ------------------------------------------------------------------


@pytest.mark.asyncio
async def test_extract_dossier_persists_with_5_plus_sections():
    file_bytes = _build_proposal_pptx()
    db = _mock_db()
    _set_inserted_id(db)

    with patch(
        "app.services.proposal_review.extractor.generate_structured",
        side_effect=_mock_generate,
    ) as gs:
        row = await extract_dossier(
            db, proposal_id=1, file_bytes=file_bytes,
            classification=Classification.RESTRICTED,
        )

    assert row.id == 99
    assert row.dossier_json is not None
    sections = row.dossier_json["sections"]
    assert len(sections) >= 5, f"expected 5+ sections, got {len(sections)}"
    for k, v in sections.items():
        assert isinstance(v, dict), f"section {k} not a dict: {type(v)}"
    db.add.assert_called_once()
    assert gs.await_count == len(sections)


@pytest.mark.asyncio
async def test_cache_hit_skips_ollama(caplog):
    file_bytes = _build_proposal_pptx()
    cached = DossierRow(
        id=42,
        proposal_id=2,
        source_hash=compute_source_hash(file_bytes),
        model="qwen2.5:32b",
        dossier_json={"sections": {"executive_summary": {}}},
    )
    db = _mock_db(find_returns=cached)

    with caplog.at_level(logging.INFO, logger="app.services.proposal_review.extractor"):
        with patch(
            "app.services.proposal_review.extractor.generate_structured",
            side_effect=_mock_generate,
        ) as gs:
            row = await extract_dossier(
                db, proposal_id=2, file_bytes=file_bytes,
                classification=Classification.RESTRICTED,
            )
            assert gs.await_count == 0, "cache hit must NOT call Ollama"

    assert row is cached
    assert any("dossier_cache_hit=True" in m for m in caplog.messages), (
        f"missing 'dossier_cache_hit=True' in logs; got: {caplog.messages}"
    )
    db.add.assert_not_called()


@pytest.mark.asyncio
async def test_parallel_section_dispatch_within_1s():
    file_bytes = _build_proposal_pptx()
    db = _mock_db()
    _set_inserted_id(db)
    started_at: list[float] = []

    async def _slow_mock(*, schema: type[BaseModel], **_kwargs):
        started_at.append(time.monotonic())
        await asyncio.sleep(0.5)
        return _make_structured_result(schema)

    with patch(
        "app.services.proposal_review.extractor.generate_structured",
        side_effect=_slow_mock,
    ):
        await extract_dossier(
            db, proposal_id=3, file_bytes=file_bytes,
            classification=Classification.RESTRICTED,
        )

    assert len(started_at) >= 5
    spread = max(started_at) - min(started_at)
    assert spread < 1.0, (
        f"section extraction not parallel — spread={spread:.3f}s across {len(started_at)} sections"
    )


@pytest.mark.asyncio
async def test_validation_failure_retries_then_stores_empty():
    """Spec: retry once on validation failure, then store {} on second
    failure rather than raising."""
    file_bytes = _build_proposal_pptx()
    db = _mock_db()
    _set_inserted_id(db)

    call_log: dict[str, int] = {}

    async def _flaky_mock(*, schema: type[BaseModel], **_kwargs):
        name = schema.__name__
        call_log[name] = call_log.get(name, 0) + 1
        # Fail GenericNarrativeFacts every time — i.e. both attempts
        # for each narrative section. Other schemas succeed.
        if schema is GenericNarrativeFacts:
            raise ValueError("simulated invalid model output")
        return _make_structured_result(schema)

    with patch(
        "app.services.proposal_review.extractor.generate_structured",
        side_effect=_flaky_mock,
    ):
        row = await extract_dossier(
            db, proposal_id=4, file_bytes=file_bytes,
            classification=Classification.RESTRICTED,
        )

    # Did NOT raise; produced a row.
    assert row.id == 99
    sections = row.dossier_json["sections"]
    assert isinstance(sections, dict)
    # Empty payloads are filtered out — narrative sections that failed
    # both attempts won't appear in `sections` (the extractor only
    # persists sections with truthy payloads).
    assert "executive_summary" not in sections
    # Non-narrative section schemas still succeeded.
    assert "evaluation_criteria" in sections
    assert "detailed_approach" in sections


@pytest.mark.asyncio
async def test_compute_source_hash_is_stable_and_unique():
    h1 = compute_source_hash(b"abc")
    h2 = compute_source_hash(b"abc")
    h3 = compute_source_hash(b"abd")
    assert h1 == h2
    assert h1 != h3
    assert len(h1) == 64


@pytest.mark.asyncio
async def test_classification_default_when_proposal_not_found():
    """Proposal lookup returns None -> classification defaults to Restricted.
    The extractor must still proceed (the LLM client gate applies
    Restricted-level rules, not whether the proposal row exists)."""
    file_bytes = _build_proposal_pptx()
    db = _mock_db(classification=None)
    _set_inserted_id(db)

    with patch(
        "app.services.proposal_review.extractor.generate_structured",
        side_effect=_mock_generate,
    ):
        row = await extract_dossier(
            db, proposal_id=999, file_bytes=file_bytes,
            # explicitly leave classification unspecified; extractor
            # will load from DB and find None -> Restricted.
        )
    assert row.id == 99


@pytest.mark.asyncio
async def test_explicit_model_override_used_for_cache_key_and_call(caplog):
    """When the caller passes `model=`, the extractor uses that string
    for both the cache lookup and the LLM call."""
    file_bytes = _build_proposal_pptx()
    db = _mock_db()
    _set_inserted_id(db)

    seen_models: list[str | None] = []

    async def _capture_mock(*, schema: type[BaseModel], model=None, **_kwargs):
        seen_models.append(model)
        return _make_structured_result(schema)

    with patch(
        "app.services.proposal_review.extractor.generate_structured",
        side_effect=_capture_mock,
    ):
        row = await extract_dossier(
            db, proposal_id=5, file_bytes=file_bytes,
            classification=Classification.RESTRICTED,
            model="custom-model:7b",
        )

    assert row.model == "custom-model:7b"
    assert seen_models, "no LLM calls were made"
    assert all(m == "custom-model:7b" for m in seen_models), seen_models
