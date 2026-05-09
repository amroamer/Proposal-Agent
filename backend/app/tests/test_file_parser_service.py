"""Unit tests for file_parser_service.

Regression-focused: production crashed with
    AttributeError: 'NoneType' object has no attribute 'text'
on POST /reviews/extract-metadata for any PPTX where at least one slide
had no speaker notes — i.e. ~every real-world deck. Root cause was
unguarded `.notes_slide.notes_text_frame.text` access while only
`has_notes_slide` was checked. These tests exercise both the common case
(real PPTX, slide with notes + slide without) and the pathological case
(has_notes_slide=True but notes_text_frame is None).
"""
import io
from unittest.mock import MagicMock, patch

from pptx import Presentation

from app.services.file_parser_service import _extract_pptx, extract_text


def _build_minimal_pptx_bytes(*, slide_two_notes: bool = False) -> bytes:
    """Build a two-slide PPTX entirely in memory.

    Slide 1 always has speaker notes. Slide 2's notes slide is left
    untouched (the failing real-world shape) unless `slide_two_notes`
    is True.
    """
    prs = Presentation()
    layout = prs.slide_layouts[1]  # "Title and Content" in default template

    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "First slide"
    body1 = next(
        (ph for ph in s1.placeholders if ph.placeholder_format.idx == 1),
        None,
    )
    if body1 is not None:
        body1.text_frame.text = "Body of first slide"
    # Touching .notes_slide.notes_text_frame.text materializes the notes.
    s1.notes_slide.notes_text_frame.text = "These are speaker notes for slide 1."

    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Second slide"
    body2 = next(
        (ph for ph in s2.placeholders if ph.placeholder_format.idx == 1),
        None,
    )
    if body2 is not None:
        body2.text_frame.text = "Body of second slide"
    if slide_two_notes:
        s2.notes_slide.notes_text_frame.text = "Notes for slide 2."
    # else: do not touch s2.notes_slide so notes are absent.

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


class TestExtractPptxNotesHandling:
    def test_two_slides_one_with_notes_one_without_no_crash(self):
        """Real-world layout: deck with mixed slides (some with notes, some
        without) must parse without raising. Pre-fix this raised
        AttributeError on the first notes-less slide."""
        data = _build_minimal_pptx_bytes(slide_two_notes=False)
        text = _extract_pptx(data)
        # Sanity: titles, bodies, and slide-1 notes all present.
        assert "First slide" in text
        assert "Second slide" in text
        assert "Body of first slide" in text
        assert "Body of second slide" in text
        assert "speaker notes for slide 1" in text.lower()
        # And we did NOT smuggle in a fake notes line for slide 2.
        assert "Notes for slide 2." not in text

    def test_extract_text_dispatch_for_pptx(self):
        """Public entry point routes .pptx correctly and returns non-empty."""
        data = _build_minimal_pptx_bytes(slide_two_notes=False)
        text, kind = extract_text("deck.pptx", data)
        assert kind == "pptx"
        assert text.strip() != ""

    def test_both_slides_with_notes_works(self):
        """Sanity: both notes paths still capture text when present."""
        data = _build_minimal_pptx_bytes(slide_two_notes=True)
        text = _extract_pptx(data)
        assert "speaker notes for slide 1" in text.lower()
        assert "notes for slide 2." in text.lower()

    def test_handles_notes_text_frame_None_regression(self):
        """Pathological state: has_notes_slide=True but notes_text_frame is None.
        This is the exact shape that crashed prod (PowerPoint-generated decks
        whose notes slides were pre-instantiated without a body placeholder).
        Pre-fix raised AttributeError; post-fix returns cleanly."""
        fake_slide = MagicMock()
        fake_slide.has_notes_slide = True
        fake_slide.notes_slide.notes_text_frame = None
        fake_slide.shapes = []  # isolate the notes path; no shapes to iterate

        fake_prs = MagicMock()
        fake_prs.slides = [fake_slide]

        # _extract_pptx imports Presentation from pptx inside the function,
        # so we patch the canonical module attribute.
        with patch("pptx.Presentation", return_value=fake_prs):
            text = _extract_pptx(b"unused-bytes-because-mock")
        assert isinstance(text, str)


class TestExtractPptxRunTextNoneSafe:
    def test_run_with_None_text_does_not_crash(self):
        """A run whose .text is None (rare but possible with corrupt XML)
        must not crash. The `or ""` defensive wrap covers this."""
        fake_run = MagicMock()
        fake_run.text = None  # crashing condition pre-defensive-wrap

        fake_para = MagicMock()
        fake_para.runs = [fake_run]

        fake_text_frame = MagicMock()
        fake_text_frame.paragraphs = [fake_para]

        fake_shape = MagicMock()
        fake_shape.has_text_frame = True
        fake_shape.has_table = False
        fake_shape.text_frame = fake_text_frame

        fake_slide = MagicMock()
        fake_slide.has_notes_slide = False
        fake_slide.shapes = [fake_shape]

        fake_prs = MagicMock()
        fake_prs.slides = [fake_slide]

        with patch("pptx.Presentation", return_value=fake_prs):
            text = _extract_pptx(b"unused")
        assert isinstance(text, str)
