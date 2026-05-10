"""Phase 1 — deterministic .pptx -> per-section slide groups.

KPMG technical proposals follow a fixed templated section order. Each
section starts on its own divider slide whose title takes the form:

    NN <separator> <Arabic title>

where NN is two digits, and the separator is an en-dash, hyphen, or
em-dash (sometimes surrounded by whitespace). Examples:

    01 – معايير التقييم
    02 - الملخص التنفيذي
    03 — فهمنا لمتطلباتكم

The splitter walks slides in order, treats every slide whose first
non-empty line matches that pattern as the start of a new section, and
maps the title onto a canonical `section_key` via section_mapping.

Slides BEFORE the first detected marker are bucketed into a synthetic
`front_matter` section.

This module never calls an LLM. It is purely deterministic.
"""
from __future__ import annotations

import io
import logging
import re
from typing import Iterable

from pydantic import BaseModel, Field

from .section_mapping import (
    FRONT_MATTER_KEY,
    arabic_label,
    english_label,
    map_title_to_key,
)

logger = logging.getLogger(__name__)


# Section marker: two ASCII digits, then any combination of whitespace
# and a hyphen / en-dash / em-dash, then the rest of the line (the
# Arabic or English title). The lookup itself is done by
# section_mapping.map_title_to_key, which is robust to tashkeel and
# alef variants. Anchored at line start so body bullets containing
# "01 – something" don't false-positive.
_SECTION_MARKER = re.compile(
    r"^\s*(\d{2})\s*[\-‐‑‒–—―]\s*(.+?)\s*$"
)


class SlideText(BaseModel):
    slide_number: int = Field(ge=1)
    text: str

    def is_blank(self) -> bool:
        return not (self.text and self.text.strip())


class SectionContent(BaseModel):
    """Slides + raw concatenated text for a single canonical section."""

    section_key: str
    title_ar: str
    title_en: str
    slides: list[SlideText] = Field(default_factory=list)
    # Convenience: the joined plain text of every slide in this section.
    # Phase-3 extraction passes this verbatim to the per-section LLM call.
    raw_text: str = ""


class ProposalSections(BaseModel):
    """Output of the splitter: keyed by canonical `section_key`."""

    sections: dict[str, SectionContent] = Field(default_factory=dict)
    # The slide index where each section starts — useful for debugging
    # and for citing slide ranges back to the user.
    section_starts: dict[str, int] = Field(default_factory=dict)
    total_slides: int = 0
    # `section_key`s the splitter expected (i.e. canonical 13) but did
    # NOT find. Logged as warnings, surfaced here for visibility.
    missing_sections: list[str] = Field(default_factory=list)
    # Section-marker slides whose title couldn't be mapped to a canonical
    # key. Slide number + raw title — surfaced so operators can extend
    # the keyword table when they encounter a deck that breaks routing.
    unknown_markers: list[dict] = Field(default_factory=list)


# ---------- public API ----------


def split_pptx(file_bytes: bytes) -> ProposalSections:
    """Split a .pptx into canonical sections. Pure function; no I/O."""
    slides = _extract_slide_texts(file_bytes)
    return split_slide_texts(slides)


def split_slide_texts(slides: list[SlideText]) -> ProposalSections:
    """Same as split_pptx but starts from already-extracted per-slide text.
    Exposed so tests can build small fixtures without round-tripping through
    python-pptx, and so the Celery extractor can re-use cached extractions."""
    out = ProposalSections(total_slides=len(slides))

    current_key: str | None = FRONT_MATTER_KEY
    out.sections[FRONT_MATTER_KEY] = SectionContent(
        section_key=FRONT_MATTER_KEY,
        title_ar="",
        title_en="Front matter",
    )
    out.section_starts[FRONT_MATTER_KEY] = 1

    for slide in slides:
        marker_title = _detect_marker(slide.text)
        if marker_title is not None:
            key = map_title_to_key(marker_title)
            if key is None:
                logger.warning(
                    "section_splitter: unrecognised marker on slide %d: %r",
                    slide.slide_number, marker_title,
                )
                out.unknown_markers.append(
                    {"slide_number": slide.slide_number, "title": marker_title}
                )
                # Treat unrecognised markers as a continuation of the
                # current section rather than starting a new one — we
                # can't route a section we don't know about, and
                # silently dropping the slides would lose evidence.
                _append_slide(out.sections[current_key], slide)
                continue

            # Drop empty front_matter so callers don't see a useless
            # section in the output. (We bucketed it as a placeholder
            # in case the deck started with cover/agenda slides.)
            if (
                current_key == FRONT_MATTER_KEY
                and not out.sections[FRONT_MATTER_KEY].slides
            ):
                out.sections.pop(FRONT_MATTER_KEY, None)
                out.section_starts.pop(FRONT_MATTER_KEY, None)

            # New section: replace if previously seen (deck quirk —
            # KPMG sometimes splits long sections across multiple
            # divider slides and we want the LATEST start).
            existing = out.sections.get(key)
            if existing is None:
                out.sections[key] = SectionContent(
                    section_key=key,
                    title_ar=arabic_label(key) or marker_title,
                    title_en=english_label(key),
                )
            out.section_starts[key] = slide.slide_number
            current_key = key
            # The marker slide itself often holds the chapter title only
            # — but capture it anyway so we don't drop any text.
            _append_slide(out.sections[key], slide)
            continue

        if current_key is None:  # defensive — should never happen
            current_key = FRONT_MATTER_KEY
            out.sections.setdefault(
                FRONT_MATTER_KEY,
                SectionContent(section_key=FRONT_MATTER_KEY, title_ar="", title_en="Front matter"),
            )
        _append_slide(out.sections[current_key], slide)

    # Compute raw_text from accumulated slides (post-walk so we don't
    # rebuild the string on every append).
    for sec in out.sections.values():
        sec.raw_text = _join_slide_text(sec.slides)

    out.missing_sections = [k for k in _expected_keys() if k not in out.sections]
    if out.missing_sections:
        logger.warning(
            "section_splitter: %d expected sections missing: %s",
            len(out.missing_sections), out.missing_sections,
        )

    return out


# ---------- helpers ----------


def _expected_keys() -> tuple[str, ...]:
    # Imported here so a future refactor can override the canonical set
    # without circular imports.
    from .section_mapping import SECTION_KEYS
    return SECTION_KEYS


def _detect_marker(slide_text: str) -> str | None:
    """If `slide_text`'s first non-empty line matches the section marker
    pattern, return the title portion. Otherwise return None."""
    if not slide_text:
        return None
    for line in slide_text.splitlines():
        line = line.strip()
        if not line:
            continue
        m = _SECTION_MARKER.match(line)
        if m:
            return (m.group(2) or "").strip()
        # Only the FIRST non-empty line counts — body slides that
        # happen to contain "01 – ..." in a bullet must not match.
        return None
    return None


def _append_slide(section: SectionContent, slide: SlideText) -> None:
    section.slides.append(slide)


def _join_slide_text(slides: Iterable[SlideText]) -> str:
    return "\n\n".join(s.text for s in slides if s.text and s.text.strip())


def _extract_slide_texts(file_bytes: bytes) -> list[SlideText]:
    """Read a .pptx and return per-slide text including table cells.

    Captures, in order:
      - text frames on every shape (titles, body, text boxes)
      - table cells (CV slides + references slides — these MUST be
        captured or team_structure / detailed_experience extraction
        loses most of its evidence)
      - speaker notes (helps with our_understanding context)

    Mirrors the public file_parser_service._extract_pptx logic but
    returns structured per-slide data instead of a flat blob — the
    splitter needs slide numbers preserved.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    out: list[SlideText] = []

    for idx, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                tf = shape.text_frame
                for para in tf.paragraphs:
                    line = "".join((run.text or "") for run in para.runs).strip()
                    if line:
                        chunks.append(line)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    cells = [c for c in cells if c]
                    if cells:
                        chunks.append(" | ".join(cells))
        # Speaker notes — both guards required: has_notes_slide can be
        # True while notes_text_frame is None on PowerPoint-generated
        # decks where notes slides were pre-instantiated empty.
        if getattr(slide, "has_notes_slide", False):
            ntf = slide.notes_slide.notes_text_frame
            if ntf is not None:
                notes = (ntf.text or "").strip()
                if notes:
                    chunks.append(f"[notes] {notes}")
        out.append(SlideText(slide_number=idx, text="\n".join(chunks)))

    return out
