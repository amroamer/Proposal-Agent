"""Extract plain text from uploaded proposal files (.pptx / .docx / .pdf)."""
import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_KINDS = ("pptx", "docx", "pdf")
# Backend-side hard cap. Keep this <= the smaller of pa_nginx.client_max_body_size
# (256M) and the outer kpmg-infra nginx client_max_body_size — otherwise the
# proxy 413s before the request ever reaches us.
MAX_BYTES = 256 * 1024 * 1024  # 256 MB


class UnsupportedFile(ValueError):
    pass


class FileTooLarge(ValueError):
    pass


def detect_kind(filename: str) -> str:
    ext = Path(filename).suffix.lower().lstrip(".")
    if ext not in SUPPORTED_KINDS:
        raise UnsupportedFile(
            f"Unsupported file type '.{ext}'. Supported: {', '.join('.' + k for k in SUPPORTED_KINDS)}"
        )
    return ext


def extract_text(filename: str, content: bytes) -> tuple[str, str]:
    """Return (extracted_text, kind). Raises UnsupportedFile / FileTooLarge / ValueError."""
    if len(content) > MAX_BYTES:
        raise FileTooLarge(f"File exceeds {MAX_BYTES // (1024*1024)} MB limit.")

    kind = detect_kind(filename)
    if kind == "pptx":
        text = _extract_pptx(content)
    elif kind == "docx":
        text = _extract_docx(content)
    elif kind == "pdf":
        text = _extract_pdf(content)
    else:
        raise UnsupportedFile(f"No extractor for {kind}")

    text = text.strip()
    if not text:
        raise ValueError(
            "No readable text found in file. The document may be image-only or empty."
        )
    return text, kind


# ---------- backends ----------

def _extract_pptx(content: bytes) -> str:
    """Pull text from every slide: shapes (incl. tables) + speaker notes."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_chunks: list[str] = [f"\n## Slide {idx}\n"]

        for shape in slide.shapes:
            # Plain text shapes (titles, body, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join((run.text or "") for run in para.runs).strip()
                    if line:
                        slide_chunks.append(line)

            # Tables — flatten row-by-row
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    cells = [c for c in cells if c]
                    if cells:
                        slide_chunks.append(" | ".join(cells))

        # Speaker notes — has_notes_slide can be True while notes_text_frame
        # is None (e.g. PowerPoint-generated decks where notes slides were
        # pre-instantiated without a body placeholder). Both guards required
        # or this raises AttributeError on every real-world deck.
        if slide.has_notes_slide:
            ntf = slide.notes_slide.notes_text_frame
            if ntf is not None:
                notes = (ntf.text or "").strip()
                if notes:
                    slide_chunks.append(f"_Speaker notes:_ {notes}")

        if len(slide_chunks) > 1:
            parts.append("\n".join(slide_chunks))

    return "\n".join(parts)


def _extract_docx(content: bytes) -> str:
    """Pull text from every paragraph + every table cell."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    parts: list[str] = []

    # Defensive `or ""` wraps: python-docx normally returns str, but corrupt
    # or unusual documents have been seen returning None for empty placeholders.
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            cells = [c for c in cells if c]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def _extract_pdf(content: bytes) -> str:
    """Per-page text via PyMuPDF (fitz)."""
    import fitz

    doc = fitz.open(stream=content, filetype="pdf")
    try:
        parts: list[str] = []
        for idx, page in enumerate(doc, start=1):
            # `or ""` guards against image-only / scanned pages that return None
            # in some PyMuPDF builds for certain encrypted or malformed PDFs.
            text = (page.get_text("text") or "").strip()
            if text:
                parts.append(f"\n## Page {idx}\n{text}")
        return "\n".join(parts)
    finally:
        doc.close()
